#!/usr/bin/env python3
"""Build the shinyreact dark Keynote theme as a .pptx.

WHY THIS IS STRUCTURED THE WAY IT IS
------------------------------------
Keynote's "Save Theme" captures **master slides**, not content slides. The first
version of this script drew everything on slides, so the saved .kth carried only
the master's background colour and nothing else - the design evaporated the moment
you started a new presentation from it.

So: every piece of design furniture goes on a *slide layout* (PowerPoint's name
for what Keynote calls a master), and every piece of text a user should be able to
edit is a real `<p:ph>` placeholder. Content slides in this deck are only examples
that demonstrate each master.

python-pptx's LayoutShapes has no add_* methods, but binding a SlideShapes to the
layout's spTree gives the full drawing API with correct image relationships.

Design authority is DESIGN.md, alongside this file. Dark mode only.

Coordinate system: the spec is written against a 1920x1080 px canvas. The deck is
960x540 pt, so **pt = px / 2** exactly. PX() does that conversion so the code can
quote spec numbers verbatim.

Run:  uv run --with python-pptx python build_theme.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.shapes.shapetree import SlideShapes
from pptx.util import Pt

HERE = Path(__file__).parent
LOGO = HERE / "shiny-react.png"
if not LOGO.exists():  # running inside the shinyreact repo
    LOGO = HERE.parent.parent / "logo" / "shiny-react.png"
OUT = HERE / "shinyreact-dark.pptx"

# ---------------------------------------------------------------- spec tokens

INK = RGBColor(0x1C, 0x1D, 0x22)
INK_DEEP = RGBColor(0x14, 0x15, 0x19)
TEXT = RGBColor(0xF2, 0xF4, 0xF8)  # 15.28:1 on INK
MUTED = RGBColor(0xA3, 0xAC, 0xBB)  # 7.35:1 on INK
CYAN = RGBColor(0x00, 0xD8, 0xFF)  # GRAPHIC ONLY
CYAN_TEXT = RGBColor(0x6F, 0xD4, 0xE8)  # 9.82:1 on INK - all text runs

C_KEYWORD = RGBColor(0xC7, 0xA0, 0xFF)
C_STRING = RGBColor(0x9F, 0xE8, 0x8D)
C_NUMBER = RGBColor(0xFF, 0xB8, 0x6C)
C_COMMENT = RGBColor(0x9A, 0xA3, 0xB2)

SERIES = [
    RGBColor(0x0E, 0x9D, 0xBA), RGBColor(0xE2, 0x62, 0x2B),
    RGBColor(0x3D, 0x8A, 0xE0), RGBColor(0x74, 0xA0, 0x12),
    RGBColor(0xC7, 0x4B, 0xD1),
]

DISPLAY = "Inter ExtraBold"
BODY_MED = "Inter Medium"
BODY = "Inter"
BOLD = "Inter Bold"
MONO = "Source Code Pro"

MARGIN = 96
FOOTER_Y = 986


def PX(px):
    return Pt(px / 2)


# ------------------------------------------------------------------- plumbing


def layout_shapes(layout):
    """Full drawing API bound to a layout's shape tree."""
    return SlideShapes(layout.shapes._spTree, layout)


def clear(layout):
    """Strip a stock layout back to an empty canvas."""
    spTree = layout.shapes._spTree
    for el in list(spTree)[2:]:  # keep nvGrpSpPr + grpSpPr
        spTree.remove(el)


def keep_only(master, keep):
    """Drop every stock layout we did not author, so the theme chooser is clean."""
    keep_parts = {lay.part for lay in keep}
    lst = master.element.get_or_add_sldLayoutIdLst()
    for el in list(lst):
        rId = el.get(qn("r:id"))
        if master.part.related_part(rId) not in keep_parts:
            lst.remove(el)
            master.part.drop_rel(rId)


def set_alpha(color_format, pct):
    srgb = color_format._xFill.find(qn("a:srgbClr"))
    srgb.append(srgb.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))}))


def no_autofit(tf):
    tf.word_wrap = True
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    for tag in ("a:normAutofit", "a:spAutoFit"):
        found = bodyPr.find(qn(tag))
        if found is not None:
            bodyPr.remove(found)


def as_placeholder(shape, ph_type, idx=None):
    """Turn a text box into a real placeholder so Keynote makes it editable."""
    nvSpPr = shape._element.nvSpPr
    cNvSpPr = nvSpPr.find(qn("p:cNvSpPr"))
    if cNvSpPr is not None and "txBox" in cNvSpPr.attrib:
        del cNvSpPr.attrib["txBox"]
    nvPr = nvSpPr.find(qn("p:nvPr"))
    ph = nvPr.makeelement(qn("p:ph"), {"type": ph_type})
    if idx is not None:
        ph.set("idx", str(idx))
    nvPr.insert(0, ph)
    return shape


def ph_defaults(shape, *, size, color, font, spacing=None, line=None, bullet=None,
                space_after=None, alpha=None):
    """Write the placeholder's own list style.

    Without this, styling lives only on the runs that happen to be there now, and
    text typed after those runs are deleted reverts to Keynote's stock look.

    `bullet="bar"` puts the cyan bar in the list style itself, so text typed into
    a fresh slide gets it. Setting it per-paragraph is not enough - this list
    style wins for anything the user types later.
    """
    txBody = shape.text_frame._txBody
    lstStyle = txBody.find(qn("a:lstStyle"))
    if lstStyle is None:
        lstStyle = txBody.makeelement(qn("a:lstStyle"), {})
        txBody.insert(1, lstStyle)
    if bullet == "bar":
        marL, indent = str(int(PX(52))), str(int(-PX(52)))
    else:
        marL, indent = "0", "0"
    lvl = lstStyle.makeelement(qn("a:lvl1pPr"), {"marL": marL, "indent": indent})
    if line:
        lnSpc = lvl.makeelement(qn("a:lnSpc"), {})
        pct = lnSpc.makeelement(qn("a:spcPct"), {"val": str(int(line * 100000))})
        lnSpc.append(pct)
        lvl.append(lnSpc)
    # schema order: lnSpc, spcBef/After, buClr, buSzPct, buFont, buChar/buNone, defRPr
    if space_after:
        spcAft = lvl.makeelement(qn("a:spcAft"), {})
        # spcPts is hundredths of a point; px -> pt is /2
        spcAft.append(spcAft.makeelement(qn("a:spcPts"),
                                         {"val": str(int(space_after / 2 * 100))}))
        lvl.append(spcAft)
    if bullet == "bar":
        buClr = lvl.makeelement(qn("a:buClr"), {})
        buClr.append(buClr.makeelement(qn("a:srgbClr"), {"val": "00D8FF"}))
        lvl.append(buClr)
        lvl.append(lvl.makeelement(qn("a:buSzPct"), {"val": "115000"}))
        lvl.append(lvl.makeelement(qn("a:buFont"), {"typeface": "Arial"}))
        lvl.append(lvl.makeelement(qn("a:buChar"), {"char": "▍"}))
    else:
        lvl.append(lvl.makeelement(qn("a:buNone"), {}))
    attrs = {"sz": str(int(size / 2 * 100))}
    if spacing:
        attrs["spc"] = str(int(spacing / 2 * 100))
    defRPr = lvl.makeelement(qn("a:defRPr"), attrs)
    fill = defRPr.makeelement(qn("a:solidFill"), {})
    clr = fill.makeelement(qn("a:srgbClr"), {"val": f"{color}"})
    if alpha is not None:
        clr.append(clr.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000))}))
    fill.append(clr)
    defRPr.append(fill)
    latin = defRPr.makeelement(qn("a:latin"), {"typeface": font})
    defRPr.append(latin)
    lvl.append(defRPr)
    lstStyle.append(lvl)


def bar_bullet(para, size):
    """Cyan bar bullet - reproduces the spec's rule without a separate shape,
    so it stays aligned to the text no matter how many lines a bullet runs to."""
    pPr = para._p.get_or_add_pPr()
    pPr.set("marL", str(int(PX(52))))
    pPr.set("indent", str(int(-PX(52))))
    buClr = pPr.makeelement(qn("a:buClr"), {})
    clr = buClr.makeelement(qn("a:srgbClr"), {"val": "00D8FF"})
    buClr.append(clr)
    pPr.append(buClr)
    pPr.append(pPr.makeelement(qn("a:buSzPct"), {"val": "115000"}))
    pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"}))
    pPr.append(pPr.makeelement(qn("a:buChar"), {"char": "▍"}))


def no_bullet(para):
    pPr = para._p.get_or_add_pPr()
    pPr.set("marL", "0")
    pPr.set("indent", "0")
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


# ------------------------------------------------------------------- drawing


def textbox(shapes, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    box = shapes.add_textbox(PX(x), PX(y), PX(w), PX(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    no_autofit(tf)
    return box, tf


def style(para, text, *, size, color, font=BODY, spacing=None,
          align=PP_ALIGN.LEFT, line=None):
    para.alignment = align
    if line:
        para.line_spacing = line
    run = para.add_run()
    run.text = text
    run.font.size = PX(size)
    run.font.color.rgb = color
    run.font.name = font
    if spacing:
        run.font._rPr.set("spc", str(int(spacing / 2 * 100)))
    return run


def rect(shapes, x, y, w, h, color, *, shape=MSO_SHAPE.RECTANGLE, alpha=None):
    sh = shapes.add_shape(shape, PX(x), PX(y), PX(w), PX(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if alpha is not None:
        set_alpha(sh.fill.fore_color, alpha)
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def hex_ring(shapes, x, y, w, h, color, width=3):
    """Pointy-top hexagon OUTLINE around the logo.

    Deliberately a stroke, not a fill. The logo PNG has its own hex baked in at
    #1C1D22, so anything filled behind it shows as a patch that does not match
    the slide. Stroking leaves the logo's hex sitting on an identical slide
    fill - the two can never disagree.

    Freeform, not MSO_SHAPE.HEXAGON: that preset is flat-top and rotating it
    90 degrees swaps the bounding box.
    """
    pts = [(x + w / 2, y), (x + w, y + h * 0.25), (x + w, y + h * 0.75),
           (x + w / 2, y + h), (x, y + h * 0.75), (x, y + h * 0.25)]
    b = shapes.build_freeform(PX(pts[0][0]), PX(pts[0][1]))
    b.add_line_segments([(PX(a), PX(c)) for a, c in pts[1:]], close=True)
    sh = b.convert_to_shape()
    sh.fill.background()
    sh.line.color.rgb = color
    sh.line.width = PX(width)
    sh.shadow.inherit = False
    return sh


def bezier(p0, p1, p2, p3, n=48):
    out = []
    for i in range(n + 1):
        t = i / n
        u = 1 - t
        out.append((
            u**3 * p0[0] + 3 * u**2 * t * p1[0] + 3 * u * t**2 * p2[0] + t**3 * p3[0],
            u**3 * p0[1] + 3 * u**2 * t * p1[1] + 3 * u * t**2 * p2[1] + t**3 * p3[1],
        ))
    return out


def swoosh(shapes, x, y, width):
    s = width / 520.0
    pts = [(x + a * s, y + b * s) for a, b in bezier((4, 34), (150, 6), (340, 6), (516, 20))]
    b = shapes.build_freeform(PX(pts[0][0]), PX(pts[0][1]))
    b.add_line_segments([(PX(a), PX(c)) for a, c in pts[1:]], close=False)
    sh = b.convert_to_shape()
    sh.fill.background()
    sh.line.color.rgb = CYAN
    sh.line.width = PX(8)
    sh.shadow.inherit = False
    sh.line._get_or_add_ln().set("cap", "rnd")
    return sh


def orbit(shapes, cx, cy, size, alpha):
    rx, ry = size * 0.92, size * 0.35
    for angle in (0, 60, 120):
        el = shapes.add_shape(MSO_SHAPE.OVAL, PX(cx - rx), PX(cy - ry),
                              PX(rx * 2), PX(ry * 2))
        el.fill.background()
        el.line.color.rgb = CYAN
        el.line.width = PX(3)
        el.rotation = angle
        el.shadow.inherit = False
        set_alpha(el.line.color, alpha)


def footer(shapes):
    _, tf = textbox(shapes, 1920 - MARGIN - 700, FOOTER_Y, 700, 40)
    style(tf.paragraphs[0], "posit::conf(2026) · @schloerke", size=30,
          color=MUTED, align=PP_ALIGN.RIGHT)


# -------------------------------------------------------------------- masters


def master_title(lay):
    sh = layout_shapes(lay)
    lay.background.fill.solid()
    lay.background.fill.fore_color.rgb = INK
    orbit(sh, 1770, 370, 550, 7)
    hex_ring(sh, 1920 - 130 - 560, 247, 560, 647, CYAN, width=4)
    sh.add_picture(str(LOGO), PX(1920 - 150 - 520), PX(270), width=PX(520))

    _, tf = textbox(sh, MARGIN, 250, 1100, 60)
    style(tf.paragraphs[0], "POSIT::CONF(2026)", size=34, color=CYAN_TEXT,
          font=BOLD, spacing=7.5)

    t, tf = textbox(sh, MARGIN, 306, 1100, 280)
    style(tf.paragraphs[0], "Title", size=128, color=TEXT, font=DISPLAY, line=0.98)
    no_bullet(tf.paragraphs[0])
    ph_defaults(t, size=128, color="F2F4F8", font=DISPLAY, line=0.98)
    as_placeholder(t, "title")

    swoosh(sh, MARGIN, 596, 520)

    b, tf = textbox(sh, MARGIN, 686, 1100, 160)
    style(tf.paragraphs[0], "Subtitle", size=46, color=MUTED, line=1.35)
    no_bullet(tf.paragraphs[0])
    ph_defaults(b, size=46, color="A3ACBB", font=BODY, line=1.35)
    as_placeholder(b, "body", 1)

    _, tf = textbox(sh, MARGIN, 940, 400, 90)
    style(tf.paragraphs[0], "posit", size=40, color=TEXT, font=DISPLAY)
    style(tf.add_paragraph(), "conf (2026)", size=26, color=MUTED, font=BODY_MED)
    rect(sh, MARGIN + 220, 918, 132, 132, TEXT)  # QR placeholder
    _, tf = textbox(sh, MARGIN + 220, 968, 132, 40)
    style(tf.paragraphs[0], "QR", size=28, color=INK, font=BOLD, align=PP_ALIGN.CENTER)

    _, tf = textbox(sh, 1920 - MARGIN - 700, 918, 700, 140)
    for i, txt in enumerate(("Barret Schloerke", "posit / Shiny Team", "@schloerke")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        style(p, txt, size=30, color=MUTED, align=PP_ALIGN.RIGHT, line=1.45)


def master_divider(lay):
    sh = layout_shapes(lay)
    lay.background.fill.solid()
    lay.background.fill.fore_color.rgb = INK_DEEP
    orbit(sh, 960, 540, 750, 10)

    n, tf = textbox(sh, MARGIN, 400, 600, 220)
    r = style(tf.paragraphs[0], "02", size=240, color=CYAN, font="Inter Black")
    set_alpha(r.font.color, 22)
    no_bullet(tf.paragraphs[0])
    # alpha must be in the list style too, or a typed section number comes out
    # solid cyan instead of the ghosted 22%
    ph_defaults(n, size=240, color="00D8FF", font="Inter Black", alpha=22)
    as_placeholder(n, "body", 1)

    t, tf = textbox(sh, MARGIN, 636, 1400, 130)
    style(tf.paragraphs[0], "Section", size=96, color=TEXT, font=DISPLAY)
    no_bullet(tf.paragraphs[0])
    ph_defaults(t, size=96, color="F2F4F8", font=DISPLAY)
    as_placeholder(t, "title")

    swoosh(sh, MARGIN, 790, 640)


def master_content(lay):
    sh = layout_shapes(lay)
    lay.background.fill.solid()
    lay.background.fill.fore_color.rgb = INK

    t, tf = textbox(sh, MARGIN, 176, 1500, 120)
    style(tf.paragraphs[0], "Title", size=84, color=TEXT, font=DISPLAY)
    no_bullet(tf.paragraphs[0])
    ph_defaults(t, size=84, color="F2F4F8", font=DISPLAY)
    as_placeholder(t, "title")

    swoosh(sh, MARGIN, 320, 440)

    b, tf = textbox(sh, MARGIN, 440, 1728, 420)
    for i in range(3):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.3
        p.space_after = PX(44)
        style(p, "Bullet", size=52, color=TEXT, font=BODY_MED)
        bar_bullet(p, 52)
    ph_defaults(b, size=52, color="F2F4F8", font=BODY_MED, line=1.3, bullet="bar",
                space_after=44)
    as_placeholder(b, "body", 1)
    footer(sh)


def master_code(lay):
    sh = layout_shapes(lay)
    lay.background.fill.solid()
    lay.background.fill.fore_color.rgb = INK

    t, tf = textbox(sh, MARGIN, 96, 1500, 110)
    style(tf.paragraphs[0], "Title", size=76, color=TEXT, font=DISPLAY)
    no_bullet(tf.paragraphs[0])
    ph_defaults(t, size=76, color="F2F4F8", font=DISPLAY)
    as_placeholder(t, "title")

    gap, pw = 48, (1920 - 2 * MARGIN - 48) // 2
    for i, x in enumerate((MARGIN, MARGIN + pw + gap)):
        panel = rect(sh, x, 250, pw, 660, INK_DEEP, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        panel.adjustments[0] = 0.06
        panel.line.color.rgb = CYAN
        panel.line.width = PX(2)
        set_alpha(panel.line.color, 18)

        f, tf = textbox(sh, x + 44, 290, pw - 88, 40)
        style(tf.paragraphs[0], "app.R", size=36, color=CYAN_TEXT, font=BOLD, spacing=4)
        no_bullet(tf.paragraphs[0])
        ph_defaults(f, size=36, color="6FD4E8", font=BOLD, spacing=4)
        as_placeholder(f, "body", 1 + i * 2)

        c, tf = textbox(sh, x + 44, 358, pw - 88, 512)
        style(tf.paragraphs[0], "code", size=34, color=TEXT, font=MONO, line=1.35)
        no_bullet(tf.paragraphs[0])
        ph_defaults(c, size=34, color="F2F4F8", font=MONO, line=1.35)
        as_placeholder(c, "body", 2 + i * 2)
    footer(sh)


def master_data(lay):
    sh = layout_shapes(lay)
    lay.background.fill.solid()
    lay.background.fill.fore_color.rgb = INK

    t, tf = textbox(sh, MARGIN, 96, 1500, 110)
    style(tf.paragraphs[0], "Title", size=76, color=TEXT, font=DISPLAY)
    no_bullet(tf.paragraphs[0])
    ph_defaults(t, size=76, color="F2F4F8", font=DISPLAY)
    as_placeholder(t, "title")

    b, tf = textbox(sh, MARGIN, 232, 1728, 60)
    style(tf.paragraphs[0], "Caption", size=30, color=MUTED, font=BODY_MED)
    no_bullet(tf.paragraphs[0])
    ph_defaults(b, size=30, color="A3ACBB", font=BODY_MED)
    as_placeholder(b, "body", 1)
    footer(sh)


MASTERS = [
    ("Title", master_title),
    ("Section divider", master_divider),
    ("Content", master_content),
    ("Code two-up", master_code),
    ("Data", master_data),
]


# ------------------------------------------------------- example content slides


UI_TSX = [
    [("const", C_KEYWORD), (" [n, setN] = ", TEXT), ("useShinyInput", CYAN_TEXT), ("(", TEXT)],
    [("  ", TEXT), ('"bins"', C_STRING), (", ", TEXT), ("30", C_NUMBER)],
    [(");", TEXT)],
    [("const", C_KEYWORD), (" plot = ", TEXT), ("useShinyOutputValue", CYAN_TEXT), ("(", TEXT)],
    [("  ", TEXT), ('"hist"', C_STRING)],
    [(");", TEXT)],
    [],
    [("return", C_KEYWORD), (" <", TEXT), ("Slider", CYAN_TEXT), (" value={n}", TEXT)],
    [("  onChange={setN} />;", TEXT)],
]

APP_PY = [
    [("# reactive computation only", C_COMMENT)],
    [("@shinyreact.reactive_output", CYAN_TEXT)],
    [("def", C_KEYWORD), (" ", TEXT), ("hist", CYAN_TEXT), ("():", TEXT)],
    [("    ", TEXT), ("return", C_KEYWORD), (" histogram(", TEXT)],
    [("        data, bins=input.bins()", TEXT)],
    [("    )", TEXT)],
]


def fill(ph, text, *, size, color, font, line=None, spacing=None):
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    no_bullet(p)
    style(p, text, size=size, color=color, font=font, line=line, spacing=spacing)


def fill_code(ph, lines, size=34):
    tf = ph.text_frame
    tf.clear()
    for i, runs in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.35
        no_bullet(p)
        for text, color in (runs or [(" ", TEXT)]):
            r = p.add_run()
            r.text = text
            r.font.size = PX(size)
            r.font.color.rgb = color
            r.font.name = MONO


def by_idx(slide, idx):
    for ph in slide.placeholders:
        if ph.element.ph_idx == idx:
            return ph
    return None


def build_examples(prs, layouts):
    title_l, div_l, content_l, code_l, data_l = layouts

    s = prs.slides.add_slide(title_l)
    fill(s.shapes.title, "React UI,\nShiny brain.", size=128, color=TEXT,
         font=DISPLAY, line=0.98)
    fill(by_idx(s, 1), "Building Shiny apps whose entire interface lives in ui.tsx",
         size=46, color=MUTED, font=BODY, line=1.35)

    s = prs.slides.add_slide(div_l)
    fill(s.shapes.title, "The hook family", size=96, color=TEXT, font=DISPLAY)
    n = by_idx(s, 1)
    fill(n, "02", size=240, color=CYAN, font="Inter Black")
    set_alpha(n.text_frame.paragraphs[0].runs[0].font.color, 22)

    s = prs.slides.add_slide(content_l)
    fill(s.shapes.title, "Pick the narrowest hook", size=84, color=TEXT, font=DISPLAY)
    tf = by_idx(s, 1).text_frame
    tf.clear()
    bullets = [
        ("A button that only pushes events uses ", "useSetShinyInput"),
        ("A card that only reads uses ", "useShinyOutputValue"),
        ("Narrow hooks make data-flow direction ", "visible at the call site"),
    ]
    for i, (plain, emph) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.3
        p.space_after = PX(44)
        style(p, plain, size=52, color=TEXT, font=BODY_MED)
        r = p.add_run()
        r.text = emph
        r.font.size, r.font.name, r.font.color.rgb = PX(52), BOLD, CYAN_TEXT
        bar_bullet(p, 52)

    s = prs.slides.add_slide(code_l)
    fill(s.shapes.title, "Two files, one app", size=76, color=TEXT, font=DISPLAY)
    fill(by_idx(s, 1), "ui.tsx", size=36, color=CYAN_TEXT, font=BOLD, spacing=4)
    fill_code(by_idx(s, 2), UI_TSX)
    fill(by_idx(s, 3), "app.py", size=36, color=CYAN_TEXT, font=BOLD, spacing=4)
    fill_code(by_idx(s, 4), APP_PY)

    s = prs.slides.add_slide(data_l)
    fill(s.shapes.title, "Bundle size by dependency", size=76, color=TEXT, font=DISPLAY)
    fill(by_idx(s, 1), "Gzipped, production build", size=30, color=MUTED, font=BODY_MED)
    sh = s.shapes
    stats = [("41 kB", "GZIPPED TOTAL"), ("0", "UI COMPONENTS SHIPPED"),
             ("2", "LANGUAGES SUPPORTED")]
    x = MARGIN
    for num, lbl in stats:
        _, tf = textbox(sh, x, 316, 560, 150)
        style(tf.paragraphs[0], num, size=132, color=CYAN_TEXT, font=DISPLAY)
        _, tf = textbox(sh, x, 470, 560, 46)
        style(tf.paragraphs[0], lbl, size=30, color=MUTED, font=BODY_MED, spacing=1.2)
        x += 620

    bars = [("react", 880, "22.4 kB"), ("react-dom", 540, "13.7 kB"),
            ("shiny-react", 196, "4.9 kB"), ("css", 64, "1.6 kB"),
            ("glue", 30, "0.8 kB")]
    axis_x, y = 566, 570
    rect(sh, axis_x, y - 10, 2, 400, MUTED, alpha=50)
    for i, (label, width, value) in enumerate(bars):
        _, tf = textbox(sh, axis_x - 430, y + 6, 400, 50)
        style(tf.paragraphs[0], label, size=30, color=MUTED, font=BODY_MED,
              align=PP_ALIGN.RIGHT)
        bar = rect(sh, axis_x + 12, y, width, 46, SERIES[i],
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        bar.adjustments[0] = 0.09
        _, tf = textbox(sh, axis_x + width + 42, y + 6, 300, 50)
        style(tf.paragraphs[0], value, size=30, color=TEXT, font=MONO)
        y += 78


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = PX(1920), PX(1080)
    master = prs.slide_masters[0]
    master.background.fill.solid()
    master.background.fill.fore_color.rgb = INK

    layouts = []
    for i, (name, fn) in enumerate(MASTERS):
        lay = master.slide_layouts[i]
        clear(lay)
        lay.element.set("preserve", "1")
        lay.element.cSld.set("name", name)
        fn(lay)
        layouts.append(lay)

    keep_only(master, layouts)
    build_examples(prs, layouts)
    prs.save(OUT)
    print(f"wrote {OUT.name}: {len(layouts)} masters, {len(prs.slides._sldIdLst)} example slides")


if __name__ == "__main__":
    build()
